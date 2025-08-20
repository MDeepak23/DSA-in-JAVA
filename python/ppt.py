from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture("./git.png", Inches(0), Inches(0), width=Inches(10), height=Inches(5.63))
prs.save("Git_and_GitHub_Session_Poster.pptx")